#!/usr/bin/env python3
# optimized_nova_rag_intelligent_formatted_v2.py
"""
Intelligent RAG service with improved response formatting (Style A - Professional).
Every response will follow the same clean, minimal format.

- Single formatting pipeline (apply_global_formatting)
- format_response_block(title, body) produces Style A output
- All nodes and LLM outputs pass through the formatter
- Added guardrails for "no information" responses
"""

import os
import json
import re
import boto3
import uuid
import numpy as np
from pinecone import Pinecone, ServerlessSpec
from pymongo import MongoClient
from dotenv import load_dotenv
from langgraph.graph import StateGraph, END
from langchain_core.messages import AIMessage, HumanMessage
from typing import Dict, Any, List, Literal, Optional, TypedDict, Annotated
import traceback
import logging
import time
from fastapi import FastAPI, Form, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
import uvicorn
from datetime import datetime
from pydantic import BaseModel
from jose import jwt, JWTError
from fastapi import Depends, Request

import whisper
import logging

import io
import pandas as pd
import docx
from pptx import Presentation
import PyPDF2
import whisper
from moviepy import AudioFileClip
from PIL import Image


UPLOAD_DIR = os.path.join(os.getcwd(), "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

CHAT_ROOT = os.path.join(os.getcwd(), "chat_store", "users")
os.makedirs(CHAT_ROOT, exist_ok=True)

# ============================================================
# Utilities & Formatting (Style A - Professional)
# ============================================================

def get_current_user_id(request: Request) -> str:
    auth = request.headers.get("Authorization")

    if not auth or not auth.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Authorization token missing")

    token = auth.split(" ")[1]

    try:
        payload = jwt.decode(
            token,
            os.getenv("JWT_SECRET"),
            algorithms=[os.getenv("JWT_ALGORITHM", "HS256")]
        )
        print("JWT PAYLOAD:", payload)
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid or expired token")

    user_id = (
        payload.get("user_id")
        or payload.get("_id")
        or payload.get("id")
        or payload.get("uid")
        or (payload.get("user", {}) or {}).get("_id")
        or (payload.get("user", {}) or {}).get("id")
    )
    

    if not user_id:
        raise HTTPException(status_code=401, detail="Invalid token payload")

    return str(user_id)

def apply_global_formatting(text: str) -> str:
    """
    Normalize and enforce the global Style A format:
    **Title** (handled in format_response_block)
    • Bullet lines
    - Clean spacing and punctuation
    """
    if not text:
        return ""

    # Normalize escaped characters
    text = text.replace("\\n", "\n").replace("\\t", " ").replace("\\r", "")
    text = text.replace('\\"', '"').replace("\\'", "'")

    # Protect bold markers by replacing them with a placeholder
    # Find all **text** patterns and replace with placeholder
    import re
    bold_pattern = r'\*\*([^*]+)\*\*'
    bold_matches = list(re.finditer(bold_pattern, text))
    
    # Replace bold markers with placeholders
    protected_text = text
    placeholders = []
    for i, match in enumerate(bold_matches):
        placeholder = f"__BOLD_{i}__"
        protected_text = protected_text.replace(match.group(0), placeholder)
        placeholders.append(match.group(1))  # Save the text inside **
    
    # Now apply formatting to the protected text
    
    # Collapse multiple spaces
    protected_text = re.sub(r"[ \t]{2,}", " ", protected_text)

    # Fix bullet-style consistency: convert -, *, • to "- "
    protected_text = re.sub(r"^\s*[-*•]\s*", "- ", protected_text, flags=re.MULTILINE)

    # Ensure spacing after punctuation (but not in numbers like "5:31")
    protected_text = re.sub(r"([.!?])([A-Za-z0-9])", r"\1 \2", protected_text)

    # Normalize multiple newlines -> at most one blank line
    protected_text = re.sub(r"\n{3,}", "\n\n", protected_text)

    # Trim each line's trailing/leading spaces
    lines = [ln.rstrip() for ln in protected_text.splitlines()]
    # Remove empty lines at start/end
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop(-1)

    # Reconstruct the text
    protected_text = "\n".join(lines).strip()
    
    # Restore the bold markers
    final_text = protected_text
    for i, original_text in enumerate(placeholders):
        placeholder = f"__BOLD_{i}__"
        final_text = final_text.replace(placeholder, f"**{original_text}**")
    
    return final_text


def format_response_block(title: str, body: str) -> str:
    """
    Produce Style A block:
    **Title**

    • item 1
    • item 2

    Closing line
    """
    # Clean the title - remove any existing asterisks and extra spaces
    title = title.strip()
    title = re.sub(r'^\**', '', title)  # Remove leading asterisks
    title = re.sub(r'\**$', '', title)  # Remove trailing asterisks
    title = title.strip()
    
    # Format the block with proper spacing
    block = f"**{title}**\n\n{body.strip()}"
    
    # Apply formatting but ensure we don't break the bold markers
    formatted = apply_global_formatting(block)
    
    # Final cleanup: ensure exactly two asterisks around title
    # Find the title line (first non-empty line)
    lines = formatted.split('\n')
    for i, line in enumerate(lines):
        line_stripped = line.strip()
        if line_stripped and not line_stripped.startswith('•'):
            # This should be the title line
            # Ensure it has exactly two asterisks at start and end
            if not (line_stripped.startswith('**') and line_stripped.endswith('**')):
                # Remove any asterisks and add proper ones
                clean_title = re.sub(r'\*+', '', line_stripped).strip()
                lines[i] = f"**{clean_title}**"
            break
    
    return '\n'.join(lines).strip()

def format_bullet_list(items: List[str]) -> str:
    if not items:
        return ""
    return "\n".join([f"• {item}" for item in items])

def format_numbered_list(items: List[str]) -> str:
    if not items:
        return ""
    return "\n".join([f"{i+1}. {item}" for i, item in enumerate(items)])

# ============================================================
# Environment & Clients
# ============================================================
load_dotenv()

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("formatted-nova-assistant")

AWS_REGION = os.getenv("AWS_REGION", "us-east-1")
BEDROCK_MODEL_ID = os.getenv("BEDROCK_CLAUDE_MODEL_ID", "amazon.nova-lite-v1:0")
BEDROCK_EMBEDDING_MODEL_ID = os.getenv("BEDROCK_EMBEDDING_MODEL_ID", "amazon.titan-embed-text-v2:0")
MONGO_URI = os.getenv("MONGO_URI", "mongodb://localhost:27017")
MONGO_DB = os.getenv("MONGO_DB", "Clinic")
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX = os.getenv("PINECONE_INDEX", "voxora-2")

# ============================================================
# Initialize external services (Bedrock, Mongo, Pinecone)
# ============================================================
def initialize_aws_clients():
    try:
        session = boto3.Session(region_name=AWS_REGION)
        bedrock = session.client("bedrock-runtime")
        logger.info("✅ AWS Bedrock client initialized")
        return bedrock
    except Exception as e:
        logger.error(f"❌ Failed to initialize AWS clients: {e}")
        return None

def initialize_mongo_client():
    try:
        mongo = MongoClient(MONGO_URI, serverSelectionTimeoutMS=5000)
        mongo.admin.command('ping')
        db = mongo[MONGO_DB]
        collections = db.list_collection_names()
        logger.info("✅ MongoDB client initialized")
        logger.info(f"📊 Database: {MONGO_DB}")
        logger.info(f"📊 Collections: {collections}")
        return mongo, db
    except Exception as e:
        logger.error(f"❌ Failed to connect to MongoDB: {e}")
        return None, None

def initialize_pinecone():
    try:
        if not PINECONE_API_KEY:
            logger.warning("⚠️ Pinecone API key not configured")
            return None, None
        
        logger.info(f"🔧 Initializing Pinecone with index: {PINECONE_INDEX}")
        pc = Pinecone(api_key=PINECONE_API_KEY)
        
        # If list_indexes() returns a different structure, adjust accordingly
        try:
            indexes = pc.list_indexes().names() if hasattr(pc.list_indexes(), "names") else pc.list_indexes()
        except Exception:
            indexes = pc.list_indexes()
            
        logger.info(f"📋 Available Pinecone indexes: {indexes}")
        
        if PINECONE_INDEX not in indexes:
            logger.warning(f"⚠️ Pinecone index '{PINECONE_INDEX}' not found in available indexes")
            return pc, None
            
        pine_index = pc.Index(PINECONE_INDEX)
        logger.info("✅ Pinecone client initialized successfully")
        return pc, pine_index
    except Exception as e:
        logger.error(f"❌ Failed to initialize Pinecone: {e}")
        logger.error(f"❌ Pinecone traceback: {traceback.format_exc()}")
        return None, None

bedrock = initialize_aws_clients()
mongo, db = initialize_mongo_client()
pc, pine_index = initialize_pinecone()

# ============================================================
# LLM & Embeddings
# ============================================================
def call_nova_model(prompt: str, max_tokens: int = 1000, temperature: float = 0.3):
    if bedrock is None:
        logger.warning("⚠️ Bedrock client not initialized")
        return None
    try:
        body = {
            "messages": [{"role": "user", "content": [{"text": prompt}]}],
            "inferenceConfig": {"maxTokens": max_tokens, "temperature": temperature}
        }
        response = bedrock.invoke_model(
            modelId=BEDROCK_MODEL_ID,
            body=json.dumps(body),
            contentType="application/json",
            accept="application/json"
        )
        response_body = json.loads(response.get("body").read())
        if "output" in response_body and "message" in response_body["output"]:
            content_list = response_body["output"]["message"].get("content", [])
            if content_list:
                reply_text = content_list[0].get("text", "")
                return reply_text
        return None
    except Exception as e:
        logger.error(f"❌ Nova API error: {e}")
        return None

def get_embedding(text: str):
    try:
        if bedrock is not None and BEDROCK_EMBEDDING_MODEL_ID:
            body = json.dumps({"inputText": text})
            response = bedrock.invoke_model(
                modelId=BEDROCK_EMBEDDING_MODEL_ID,
                body=body,
                contentType="application/json",
                accept="application/json"
            )
            response_body = json.loads(response.get('body').read())
            embedding = response_body.get('embedding')
            if embedding:
                return embedding
    except Exception:
        pass
    # Fallback deterministic pseudo-embedding
    import hashlib
    hash_obj = hashlib.md5(text.encode())
    seed = int(hash_obj.hexdigest()[:8], 16)
    np.random.seed(seed)
    return np.random.normal(0, 1, 1024).tolist()


logger.info("🔊 Loading Whisper model...")
WHISPER_MODEL = whisper.load_model("base")
logger.info("✅ Whisper model loaded")

def extract_text_from_any_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    # ---------- PDF ----------
    if ext == ".pdf":
        reader = PyPDF2.PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    # ---------- DOCX ----------
    if ext == ".docx":
        d = docx.Document(file_path)
        text = [p.text for p in d.paragraphs if p.text.strip()]
        for table in d.tables:
            for row in table.rows:
                text.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(text)

    # ---------- PPTX ----------
    if ext == ".pptx":
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    # ---------- CSV ----------
    if ext == ".csv":
        return pd.read_csv(file_path).to_string(index=False)

    # ---------- XLSX ----------
    if ext == ".xlsx":
        return pd.read_excel(file_path).to_string(index=False)

    # ---------- TEXT / CODE ----------
    if ext in [".txt", ".py", ".md"]:
        with open(file_path, encoding="utf-8", errors="ignore") as f:
            return f.read()

    # ---------- AUDIO ----------
    if ext in [".mp3", ".wav", ".m4a"]:
        return WHISPER_MODEL.transcribe(file_path)["text"]

    # ---------- VIDEO ----------
    # ---------- VIDEO ----------
    if ext in [".mp4", ".mov", ".mkv"]:
      wav_path = file_path + ".wav"

    try:
        clip = AudioFileClip(file_path)
        clip.write_audiofile(
            wav_path,
            codec="pcm_s16le",
            fps=16000,
            logger=None
        )
        clip.close()
    except Exception as e:
        raise RuntimeError(f"Audio extraction failed: {e}")

    result = WHISPER_MODEL.transcribe(
        wav_path,
        fp16=False
    )

    return result.get("text", "")


    # ---------- IMAGE (OCR) ----------
    if ext in [".jpg", ".jpeg", ".png"]:
        with open(file_path, "rb") as img:
            textract = boto3.client("textract", region_name=AWS_REGION)
            result = textract.detect_document_text(
                Document={"Bytes": img.read()}
            )
        return "\n".join(
            b["Text"] for b in result["Blocks"]
            if b["BlockType"] == "LINE"
        )

    raise ValueError(f"Unsupported file type: {ext}")



# ============================================================
# Query analysis & Mongo interaction
# ============================================================
def analyze_query_intent(query: str) -> Dict[str, Any]:
    collection_schemas = {}
    if db is not None:
        for collection_name in ['clinic', 'doctors', 'appointments', 'slots', 'notices', 'slotexception']:
            if collection_name in db.list_collection_names():
                try:
                    sample = db[collection_name].find_one({}, {"_id": 0})
                    if sample:
                        collection_schemas[collection_name] = list(sample.keys())
                except Exception:
                    continue

    prompt = f"""Analyze this user query and determine how to query the MongoDB database.

User Query: "{query}"

Available collections and their fields:
{json.dumps(collection_schemas, indent=2)}

Return ONLY valid JSON with keys:
"collection" (one of clinic, doctors, appointments, slots, notices, slotexception),
"fields" (list or null),
"filters" (MongoDB dict),
"explanation" (brief),
"query_type" (list_all|search_specific|find_by_name|find_by_date|other)
"""
    try:
        response = call_nova_model(prompt, max_tokens=400, temperature=0.1)
        if response:
            json_match = re.search(r'\{.*\}', response, re.DOTALL)
            if json_match:
                result = json.loads(json_match.group())

                ALLOWED_COLLECTIONS = {
                    "doctors", "clinic", "appointments", "slots", "notices", "slotexception"
                }

                collection = result.get("collection")

                if collection not in ALLOWED_COLLECTIONS:
                    ql = query.lower()
                    if "doctor" in ql:
                        result["collection"] = "doctors"
                    elif "clinic" in ql:
                        result["collection"] = "clinic"
                    elif "appointment" in ql:
                        result["collection"] = "appointments"
                    elif "slot" in ql:
                        result["collection"] = "slots"
                    elif "notice" in ql:
                        result["collection"] = "notices"
                    elif "exception" in ql or "holiday" in ql:
                        result["collection"] = "slotexception"

                logger.info(f"🤖 Query analysis: {result}")
                return result

    except Exception as e:
        logger.error(f"❌ Query analysis failed: {e}")

    return analyze_query_fallback(query)

def analyze_query_fallback(query: str) -> Dict[str, Any]:
    ql = query.lower()
    collection = None
    filters = {}
    if "doctor" in ql:
        collection = "doctors"
    elif "clinic" in ql:
        collection = "clinic"
    elif "appointment" in ql:
        collection = "appointments"
    elif "slot" in ql:
        collection = "slots"
        days = ["monday","tuesday","wednesday","thursday","friday","saturday","sunday"]
        for day in days:
            if day in ql:
                filters = {"dayOfWeek": {"$regex": day, "$options": "i"}}
                break
    elif "notice" in ql:
        collection = "notices"
    elif "exception" in ql or "holiday" in ql:
        collection = "slotexception"

    return {
        "collection": collection,
        "fields": None,
        "filters": filters,
        "explanation": "Fallback analysis",
        "query_type": "search_specific" if filters else "list_all"
    }

def intelligent_mongo_query(query: str) -> List[Dict]:
    if db is None:
        return []

    analysis = analyze_query_intent(query)
    collection_name = analysis.get("collection")
    if not collection_name:
      logger.info("⚠️ No schema detected → skip Mongo")
      return []

    filters = analysis.get("filters", {})
    fields = analysis.get("fields")

    if not collection_name or collection_name not in db.list_collection_names():
        logger.warning(f"⚠️ Invalid collection requested: {collection_name}")
        return []

    logger.info(f"🔍 Querying {collection_name} with filters: {filters}")
    try:
        projection = None
        if fields and isinstance(fields, list):
            projection = {field: 1 for field in fields if field}
            projection["_id"] = 0

        cursor = db[collection_name].find(filters, projection) if projection else db[collection_name].find(filters)
        results = list(cursor.limit(50))
        logger.info(f"📄 Found {len(results)} documents")
        return results
    except Exception as e:
        logger.error(f"❌ MongoDB query error: {e}")
        return []

# ============================================================
# Formatting functions for collections (use format_response_block)
# ============================================================

def format_doctors_response(results: List[Dict], query: str) -> str:
    if not results:
        return format_response_block("Doctor Information", "No matching doctors were found in the database.")

    # Check if query is asking for a specific doctor by name
    q_lower = query.lower()
    specific_doctor = None
    for doc in results:
        name = doc.get('name', '').lower()
        if name and name in q_lower:
            specific_doctor = doc
            break
    
    if specific_doctor or len(results) == 1:
        d = specific_doctor or results[0]
        name = d.get('name', 'Unknown Doctor')
        specialization = d.get('specialization') or d.get('specialty') or d.get('expertise') or 'General Practice'
        email = d.get('email', 'Not available')
        phone = d.get('phone') or d.get('mobile') or d.get('contactNumber', 'Not available')
        experience = d.get('experience', '')
        qualification = d.get('qualification', '')
        clinic = d.get('clinic') or d.get('clinicName', '')

        lines = [
            f"• Name: {name}",
            f"• Specialization: {specialization}"
        ]
        
        if email and email.lower() not in ['none', 'n/a', 'not available', '']:
            lines.append(f"• Email: {email}")
        
        if phone and phone.lower() not in ['none', 'n/a', 'not available', '']:
            lines.append(f"• Phone: {phone}")
        
        if clinic:
            lines.append(f"• Clinic: {clinic}")
        
        if experience:
            lines.append(f"• Experience: {experience}")
        
        if qualification:
            lines.append(f"• Qualification: {qualification}")

        body = "\n".join(lines) + "\n\nIf you'd like more details, feel free to ask."
        return format_response_block("Doctor Information", body)

    # Multiple doctors
    items = []
    for i, doc in enumerate(results[:15], 1):
        name = doc.get('name', f'Doctor {i}')
        spec = doc.get('specialization') or doc.get('specialty') or ''
        if spec:
            items.append(f"{i}. {name} — {spec}")
        else:
            items.append(f"{i}. {name}")
    
    body = format_bullet_list(items) + f"\n\nTotal: {len(results)} doctors found.\nTo get detailed information about any doctor, please ask for them by name."
    return format_response_block("Doctor List", body)

def format_clinics_response(results: List[Dict], query: str) -> str:
    if not results:
        return format_response_block("Clinics", "No clinics were found in the database.")
    
    # Check if query is asking for a specific clinic
    q_lower = query.lower()
    specific_clinic = None
    for clinic in results:
        name = (clinic.get('clinicName') or clinic.get('name', '')).lower()
        if name and name in q_lower:
            specific_clinic = clinic
            break
    
    if specific_clinic or len(results) == 1:
        c = specific_clinic or results[0]
        name = c.get('clinicName') or c.get('name', 'Unknown Clinic')
        address = c.get('address', 'Address not specified')
        phone = c.get('phone') or c.get('contactNumber') or c.get('mobile', 'Not available')
        email = c.get('email', 'Not available')
        specialty = c.get('specialty') or c.get('specialization', 'General Practice')
        
        lines = [
            f"• Name: {name}",
            f"• Address: {address}",
            f"• Phone: {phone}"
        ]
        
        if email and email.lower() not in ['none', 'n/a', 'not available', '']:
            lines.append(f"• Email: {email}")
        
        if specialty and specialty.lower() not in ['none', 'n/a', '']:
            lines.append(f"• Specialty: {specialty}")
        
        body = "\n".join(lines) + "\n\nFor more details or to book an appointment, please contact the clinic directly."
        return format_response_block("Clinic Information", body)
    
    # Multiple clinics
    items = []
    for i, c in enumerate(results[:20], 1):
        name = c.get('clinicName') or c.get('name', f'Clinic {i}')
        address = c.get('address', '')
        phone = c.get('phone') or c.get('contactNumber') or c.get('mobile', '')
        
        entry = f"{i}. {name}"
        if address:
            entry += f"\n   • Address: {address}"
        if phone and phone.lower() not in ['none', 'n/a', '']:
            entry += f"\n   • Phone: {phone}"
        
        items.append(entry)
    
    body = "\n\n".join(items) + f"\n\nTotal: {len(results)} clinics found.\nFor detailed information about a specific clinic, mention its name."
    return format_response_block("Clinics", body)

def format_slots_response(results: List[Dict], query: str) -> str:
    if not results:
        return format_response_block("Available Slots", "No available slots matched your request.")
    
    # Group slots by day if showing multiple days
    slots_by_day = {}
    for slot in results[:30]:
        day = slot.get('dayOfWeek', 'Unknown Day')
        if day not in slots_by_day:
            slots_by_day[day] = []
        slots_by_day[day].append(slot)
    
    items = []
    for day, day_slots in sorted(slots_by_day.items()):
        # Sort slots by start time
        sorted_slots = sorted(day_slots, 
                            key=lambda x: x.get('startTime') or x.get('start', ''))
        
        day_header = f"**{day}**"
        day_items = []
        for slot in sorted_slots:
            start = slot.get('startTime', slot.get('start', 'Unknown'))
            end = slot.get('endTime', slot.get('end', 'Unknown'))
            maxp = slot.get('maximumPatients') or slot.get('maxPatients', '')
            
            slot_entry = f"• {start} - {end}"
            if maxp:
                slot_entry += f" (Max: {maxp} patients)"
            
            doctor = slot.get('doctorName') or slot.get('doctor', '')
            if doctor:
                slot_entry += f" — Dr. {doctor}"
            
            day_items.append(slot_entry)
        
        if day_items:
            items.append(day_header + "\n" + "\n".join(day_items))
    
    body = "\n\n".join(items) + f"\n\nTotal available slots: {len(results)}"
    if len(slots_by_day) == 1:
        day_name = list(slots_by_day.keys())[0]
        body += f"\nTo book a {day_name} slot, please provide the exact time and your details."
    else:
        body += "\nTo book a slot, please specify the day, time, and provide your details."
    
    return format_response_block("Available Appointment Slots", body)

def format_notices_response(results: List[Dict], query: str) -> str:
    if not results:
        return format_response_block("Notices", "No notices found at this time.")
    
    # Sort by date if available
    sorted_results = sorted(results, 
                          key=lambda x: x.get('createdAt', '') or x.get('date', '') or x.get('timestamp', ''),
                          reverse=True)
    
    items = []
    for i, n in enumerate(sorted_results[:10], 1):
        title = n.get('title', f'Notice {i}')
        msg = n.get('message') or n.get('content') or n.get('description', '')
        created = n.get('createdAt') or n.get('date') or n.get('timestamp', '')
        
        # Format date nicely
        if created:
            try:
                # Try to parse and format date
                if isinstance(created, str):
                    # Remove milliseconds if present
                    created = created.split('.')[0]
                    date_obj = datetime.fromisoformat(created.replace('Z', '+00:00'))
                    created = date_obj.strftime("%Y-%m-%d %H:%M")
            except:
                # If parsing fails, use as-is
                pass
        
        entry = f"**{title}**"
        if msg:
            # Clean up the message
            msg = msg.strip()
            if msg.endswith('.'):
                entry += f"\n  {msg}"
            else:
                entry += f"\n  {msg}."
        if created:
            entry += f"\n  • Posted: {created}"
        
        items.append(entry)
    
    body = "\n\n".join(items) + f"\n\nTotal notices: {len(results)}"
    return format_response_block("Latest Notices", body)

def format_appointments_response(results: List[Dict], query: str) -> str:
    if not results:
        return format_response_block("Appointments", "No appointments matched your query.")
    if len(results) == 1:
        a = results[0]
        patient = a.get('patientName', 'Unknown')
        doctor = a.get('doctorName', 'Unknown')
        date = a.get('date', 'Unknown')
        timev = a.get('time') or a.get('assignedTime', '')
        lines = [f"• Patient: {patient}", f"• Doctor: {doctor}", f"• Date: {date}"]
        if timev:
            lines.append(f"• Time: {timev}")
        if a.get('status'):
            lines.append(f"• Status: {a.get('status')}")
        body = "\n".join(lines)
        return format_response_block("Appointment Details", body)
    else:
        items = []
        for i, a in enumerate(results[:10], 1):
            patient = a.get('patientName', f'Patient {i}')
            doctor = a.get('doctorName', 'Unknown')
            date = a.get('date', '')
            entry = f"{i}. {patient} — {doctor}"
            if date:
                entry += f" • {date}"
            items.append(entry)
        body = format_bullet_list(items) + f"\n\nTotal appointments found: {len(results)}"
        return format_response_block("Appointments", body)

def format_exceptions_response(results: List[Dict], query: str) -> str:
    if not results:
        return format_response_block("Schedule Exceptions", "No schedule exceptions found.")
    items = []
    for i, exc in enumerate(results[:10], 1):
        reason = exc.get('reason', 'Unknown reason')
        date = exc.get('date', 'Unknown date')
        entry = f"{i}. {reason} • {date}"
        items.append(entry)
    body = format_bullet_list(items) + f"\n\nTotal exceptions found: {len(results)}"
    return format_response_block("Schedule Exceptions", body)

def format_general_response(results: List[Dict], query: str, analysis: Dict) -> str:
    # Use LLM to produce a concise, structured response, then format
    results_text = json.dumps(results[:5], indent=2, default=str)
    prompt = f"""Based on this user query and database results, provide a concise, professional response.

User Query: "{query}"
Query Analysis: {json.dumps(analysis, indent=2)}
Database Results ({len(results)} items, first 5 shown):
{results_text}

Provide a short title and concise bullet points. End with an invitation for follow-up.
Return only the response text (no extra JSON)."""
    try:
        response = call_nova_model(prompt, max_tokens=400, temperature=0.2)
        if response:
            # LLM may return multiple paragraphs; rely on format_response_block around a short title
            title = "Results"
            body = response.strip()
            return format_response_block(title, body)
    except Exception as e:
        logger.warning(f"⚠️ LLM formatting failed: {e}")
    # Fallback
    body = f"I found {len(results)} results. Showing the first few:\n\n{results_text}"
    return format_response_block("Results", body)

# ============================================================
# Pinecone helpers
# ============================================================
def query_pinecone_intelligent(query: str, user_id: str = "default_user") -> List[str]:
    if pine_index is None:
        return []
    try:
        embedding = get_embedding(query)
        result = pine_index.query(
            vector=embedding,
            top_k=5,
            include_metadata=True,
            filter={"user_id": {"$eq": user_id}}
        )
        texts = []
        for match in result.get("matches", []):
            metadata = match.get("metadata", {})
            if metadata and "text" in metadata:
                texts.append(metadata["text"])
        return texts
    except Exception as e:
        logger.error(f"❌ Pinecone query error: {e}")
        return []

# ============================================================
# Guardrail Functions
# ============================================================

def normalize_text(text: str) -> str:
    """Normalize text for comparison."""
    return re.sub(r'\s+', ' ', re.sub(r'[^\w\s]', '', text.lower())).strip()

SCHEMA_KEYWORDS = {
    "doctor", "doctors", "clinic", "clinics",
    "appointment", "appointments",
    "slot", "slots",
    "notice", "notices",
    "exception", "holiday"
}

def is_schema_query(text: str) -> bool:
    """Check if query is about database schema."""
    normalized = normalize_text(text)
    return any(word in normalized for word in SCHEMA_KEYWORDS)

def is_knowledge_query(text: str) -> bool:
    """Check if this is asking for factual information we might not have."""
    text_lower = text.lower()
    
    # Questions asking for definitions, full forms, explanations
    question_patterns = [
        r'what (is|are) (the )?(full form of|meaning of|definition of|information about)',
        r'explain',
        r'describe',
        r'tell me about',
        r'who (is|are)',
        r'when (is|was|did)',
        r'where (is|are|was)',
        r'why (is|are|does)',
        r'how (does|do|is|are)'
    ]
    
    for pattern in question_patterns:
        if re.search(pattern, text_lower):
            return True
    
    # Specific factual queries
    factual_keywords = [
        'what is', 'what are', 'full form', 'meaning of', 
        'definition', 'explain', 'describe', 'information about',
        'tell me about'
    ]
    
    return any(keyword in text_lower for keyword in factual_keywords)

def format_pinecone_response(texts: List[str], query: str) -> str:
    """Format Pinecone response with guardrail for no information."""
    if not texts:
        # Check if this looks like a factual/knowledge question
        if is_knowledge_query(query):
            return format_response_block(
                "Answer", 
                "I don't have information about that in my knowledge base. This information is not available in the provided documents."
            )
        else:
            return format_response_block("Document Search", "No relevant information found in the uploaded documents.")
    
    combined = " ".join(texts[:3])[:1500]
    
    # Modified prompt to handle "no info" cases
    prompt = f"""User asked: "{query}"
    
Relevant document content:
{combined}

Based ONLY on the document content above:
1. If the exact answer to the question is in the text, provide it directly
2. If the information is NOT in the text, clearly state "I don't have information about that in my knowledge base"
3. If multiple pieces of relevant info exist, provide only the most relevant
4. Keep the answer concise and focused on answering exactly what was asked

Answer:"""
    
    try:
        response = call_nova_model(prompt, max_tokens=200, temperature=0.1)
        if response:
            response = response.strip()
            
            # Check if response indicates no information
            no_info_phrases = [
                "i don't have information",
                "information is not available",
                "not found in the document",
                "not available in the provided",
                "no information about that",
                "not mentioned in the text",
                "based on the document content",
                "the document does not mention"
            ]
            
            if any(phrase in response.lower() for phrase in no_info_phrases):
                return format_response_block(
                    "Answer", 
                    "I don't have information about that in my knowledge base. This information is not available in the provided documents."
                )
            
            # Extract just the answer part
            patterns = [
                r'Answer:\s*(.*)',
                r'The answer is:\s*(.*)',
                r'Based on.*?:\s*(.*)',
            ]
            
            for pattern in patterns:
                match = re.search(pattern, response, re.IGNORECASE | re.DOTALL)
                if match:
                    response = match.group(1).strip()
                    break
            
            # Clean up the response
            response = re.sub(r'\s+', ' ', response).strip()
            
            # Ensure it ends with proper punctuation
            if response and not response.endswith(('.', '!', '?')):
                response = response + '.'
            
            return format_response_block("Answer", response)
    except Exception as e:
        logger.warning(f"⚠️ LLM formatting failed for pinecone: {e}")
    
    # Fallback - just return the most relevant snippet
    if texts:
        # Try to find the most relevant sentence containing the query terms
        query_terms = query.lower().split()
        best_match = texts[0][:300]  # Take first 300 chars of first result
        
        for text in texts:
            text_lower = text.lower()
            matches = sum(1 for term in query_terms if term in text_lower)
            if matches > 0:
                # Extract a sentence around the match
                sentences = re.split(r'(?<=[.!?])\s+', text)
                for sentence in sentences:
                    sentence_lower = sentence.lower()
                    if any(term in sentence_lower for term in query_terms):
                        best_match = sentence[:200]
                        break
        
        if is_knowledge_query(query):
            return format_response_block(
                "Answer", 
                "I don't have information about that in my knowledge base. This information is not available in the provided documents."
            )
        
        return format_response_block("Answer", best_match + "...")
    
    # Final fallback
    if is_knowledge_query(query):
        return format_response_block(
            "Answer", 
            "I don't have information about that in my knowledge base. This information is not available in the provided documents."
        )
    
    return format_response_block("Document Search", "No relevant information found in the uploaded documents.")

# ============================================================
# Router & Nodes (LangGraph-like flow)
# ============================================================

class GraphState(TypedDict):
    messages: Annotated[list, "messages"]
    user_id: str
    error: str
    route: str

def router(state: GraphState) -> Dict[str, Any]:
    messages = state["messages"]
    if not messages:
        return {"route": "pinecone"}
    
    query_obj = messages[-1]
    qtext = query_obj.content.strip() if hasattr(query_obj, "content") else str(query_obj).strip()
    logger.info(f"🎯 Routing query: '{qtext}'")
    ql = qtext.lower().strip()
    
    normalized = normalize_text(qtext)

    # Exact greetings
    exact_greetings = {
        "hi", "hello", "hey", "hii", "helloo",
        "greetings", "good morning", "good afternoon", "good evening"
    }

    if normalized in exact_greetings:
        logger.info("🔄 Routing to: greeting")
        return {"route": "greeting"}

    # System info queries
    info_keywords = ["what can you do", "what do you know", "help", "what information", "what are you", "what can you help", "capabilities", "features"]
    for kw in info_keywords:
        if kw in ql:
            logger.info("🔄 Routing to: system_info")
            return {"route": "system_info"}

    # Document-related queries
    doc_keywords = ["document", "pdf", "upload", "file", "chapter", "textbook", "study", "material"]
    for kw in doc_keywords:
        if kw in ql:
            logger.info("🔄 Routing to: pinecone")
            return {"route": "pinecone"}

    # Schema-based DB queries → Mongo
    if is_schema_query(normalized):
        logger.info("🔄 Routing to: mongo")
        return {"route": "mongo"}
    
    # Specific handling for slot queries
    slot_patterns = ["slot", "available", "appointment", "booking", "schedule", "time", "monday", "tuesday", 
                    "wednesday", "thursday", "friday", "saturday", "sunday", "morning", "afternoon", "evening"]
    
    if any(pattern in ql for pattern in slot_patterns) and ("when" in ql or "available" in ql or "slot" in ql):
        logger.info("🔄 Routing to: mongo (slot query)")
        return {"route": "mongo"}

    # Check if this is a factual/knowledge query
    if is_knowledge_query(qtext):
        logger.info("🔄 Routing to: pinecone (knowledge query)")
        return {"route": "pinecone"}

    # Everything else → Pinecone
    logger.info("🔄 Routing to: pinecone (fallback)")
    return {"route": "pinecone"}


def route_decision(router_output: dict) -> str:
    return router_output["route"]


def greeting_node(state: GraphState) -> Dict[str, Any]:
    messages = state["messages"]
    # Style A short, predefined message
    title = "Welcome"
    body = (
        "• I can help with: Searching doctors, clinics, appointments\n"
        "• Checking time slots, notices, exceptions\n"
        "• Reading and answering from uploaded PDF documents\n\n"
        "If you'd like to proceed, tell me what you need."
    )
    formatted = format_response_block(title, body)
    return {"messages": messages + [AIMessage(content=formatted)]}

def system_info_node(state: GraphState) -> Dict[str, Any]:
    messages = state["messages"]
    title = "What I Can Help With"
    body_lines = [
        "• Doctors: Information about physicians and their specializations",
        "• Clinics: Names, addresses, phone numbers, specialties",
        "• Appointments: Patient appointments, schedules, status",
        "• Time Slots: Available booking slots with timing and capacity",
        "• Notices: Important announcements and updates",
        "• Exceptions: Schedule changes, holidays, cancellations",
        "",
        "Examples:",
        "• \"Show me all doctors\"",
        "• \"What clinics are available?\"",
        "• \"When are Monday slots available?\"",
        "• \"Do you have any notices?\"",
        "• \"Tell me about doctor Anindya\""
    ]
    body = "\n".join(body_lines)
    formatted = format_response_block(title, body)
    return {"messages": messages + [AIMessage(content=formatted)]}

def mongo_query_node(state: GraphState) -> Dict[str, Any]:
    messages = state["messages"]
    
    # Check if we should actually run this node
    route = state.get("route", "")
    if route != "mongo":
        # Even if not explicitly routed, check if the query is schema-related
        query_obj = messages[-1]
        qtext = query_obj.content if hasattr(query_obj, "content") else str(query_obj)
        if not is_schema_query(normalize_text(qtext)):
            logger.info("⏭️ Skipping Mongo node (not a schema query)")
            return state
    
    query_obj = messages[-1]
    qtext = query_obj.content if hasattr(query_obj, "content") else str(query_obj)
    try:
        results = intelligent_mongo_query(qtext)
        analysis = analyze_query_intent(qtext)
        # Choose a formatter based on collection
        collection = analysis.get("collection", "")
        if collection == "doctors":
            response = format_doctors_response(results, qtext)
        elif collection == "clinic":
            response = format_clinics_response(results, qtext)
        elif collection == "slots":
            response = format_slots_response(results, qtext)
        elif collection == "notices":
            response = format_notices_response(results, qtext)
        elif collection == "appointments":
            response = format_appointments_response(results, qtext)
        elif collection == "slotexception":
            response = format_exceptions_response(results, qtext)
        else:
            response = format_general_response(results, qtext, analysis)
        return {"messages": messages + [AIMessage(content=response)]}
    except Exception as e:
        logger.error(f"❌ MongoDB query error: {e}")
        error_msg = format_response_block("Error", "I encountered an error while searching the database. Please try again or rephrase your question.")
        return {"messages": messages + [AIMessage(content=error_msg)], "error": str(e)}

def pinecone_query_node(state: GraphState) -> Dict[str, Any]:
    messages = state["messages"]
    query_obj = messages[-1]
    qtext = query_obj.content if hasattr(query_obj, "content") else str(query_obj)
    user_id = state.get("user_id", "default_user")
    try:
        texts = query_pinecone_intelligent(qtext, user_id)
        response = format_pinecone_response(texts, qtext)
        return {"messages": messages + [AIMessage(content=response)]}
    except Exception as e:
        logger.error(f"❌ Pinecone query error: {e}")
        error_msg = format_response_block("Error", "I encountered an error while searching your documents. Please make sure you've uploaded PDF files first.")
        return {"messages": messages + [AIMessage(content=error_msg)], "error": str(e)}

# Build workflow
workflow = StateGraph(GraphState)
workflow.add_node("router", router)
workflow.add_node("mongo", mongo_query_node)
workflow.add_node("pinecone", pinecone_query_node)
workflow.add_node("greeting", greeting_node)
workflow.add_node("system_info", system_info_node)
workflow.set_entry_point("router")

workflow.add_conditional_edges(
    "router",
    route_decision,
    {
        "mongo": "mongo",
        "pinecone": "pinecone",
        "greeting": "greeting",
        "system_info": "system_info"
    }
)

workflow.add_edge("mongo", END)
workflow.add_edge("pinecone", END)
workflow.add_edge("greeting", END)
workflow.add_edge("system_info", END)

graph = workflow.compile()
logger.info("✅ Enhanced workflow compiled successfully")

# ============================================================
# Session-based chat storage (React Query compatible)
# ============================================================
chat_sessions = {}  # In-memory session storage

def get_session_chat(session_id: str) -> List[Dict]:
    return chat_sessions.get(session_id, [])

def add_message_to_session(session_id: str, message: Dict):
    if session_id not in chat_sessions:
        chat_sessions[session_id] = []
    chat_sessions[session_id].append(message)

def clear_session_chat(session_id: str):
    if session_id in chat_sessions:
        del chat_sessions[session_id]

def clear_all_sessions():
    """Clear all chat sessions"""
    global chat_sessions
    chat_sessions = {}

# ============================================================
# FastAPI Endpoints
# ============================================================
app = FastAPI(
    title="Hybrid RAG Assistant",
    description="AI Assistant with enhanced response formatting (Style A)",
    version="4.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:8081"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Define request models for better validation
class ChatRequest(BaseModel):
    query: str
    session_id: str

class UploadRequest(BaseModel):
    user_id: str = "default_user"

class DeleteFileRequest(BaseModel):
    filename: str

@app.post("/auth/logout")
async def logout_user(request: Request):
    """Logout user and clear their chat session"""
    try:
        user_id = get_current_user_id(request)
        session_id = f"session_{user_id}"
        clear_session_chat(session_id)
        return {"success": True, "message": "Logged out successfully"}
    except Exception as e:
        return {"success": True, "message": "Logged out"}

@app.get("/chat/session")
async def get_chat_session(request: Request):
    """Get current session chat messages"""
    user_id = get_current_user_id(request)
    session_id = f"session_{user_id}"
    messages = get_session_chat(session_id)
    return {"success": True, "messages": messages}

@app.post("/chat/message")
async def send_chat_message(
    payload: ChatRequest,
    request: Request,
    user_id: str = Depends(get_current_user_id)
):
    """Send message and get AI response"""
    query = payload.query.strip()
    session_id = f"session_{user_id}"

    if not query:
        raise HTTPException(400, "Query cannot be empty")

    # Add user message to session
    user_message = {
        "id": str(uuid.uuid4()),
        "text": query,
        "isUser": True,
        "timestamp": datetime.now().isoformat()
    }
    add_message_to_session(session_id, user_message)

    # Process with RAG - First get the route
    initial_state = GraphState(
        messages=[HumanMessage(content=query)],
        user_id=user_id,
        error=""
    )
    
    # Get the route decision
    router_output = router(initial_state)
    
    # Create new state with route information
    state_with_route = GraphState(
        messages=initial_state["messages"],
        user_id=initial_state["user_id"],
        error="",
        route=router_output.get("route", "pinecone")  # Add route to state
    )
    
    # Process through the graph
    result = graph.invoke(state_with_route)
    ai_text = apply_global_formatting(result["messages"][-1].content)

    # Add AI response to session
    ai_message = {
        "id": str(uuid.uuid4()),
        "text": ai_text,
        "isUser": False,
        "timestamp": datetime.now().isoformat()
    }
    add_message_to_session(session_id, ai_message)

    return {
        "success": True,
        "message": ai_message,
        "response": ai_text
    }

@app.post("/api/upload-pdf")
async def api_upload_pdf(
    request: Request,
    file: UploadFile = File(...),
    user_id: str = Depends(get_current_user_id)
):
    try:
        # ---- user folder ----
        user_dir = os.path.join(UPLOAD_DIR, user_id)
        os.makedirs(user_dir, exist_ok=True)

        file_path = os.path.join(user_dir, file.filename)

        # ---- save file ----
        content = await file.read()
        with open(file_path, "wb") as f:
            f.write(content)

        # ---- extract text (ANY FILE) ----
        try:
            full_text = extract_text_from_any_file(file_path)
        except Exception as e:
            logger.error(f"Extraction failed: {e}")
            full_text = ""

        full_text = full_text.strip()
        if not full_text:
            full_text = f"Document: {file.filename}\nNo readable text found."
            logger.warning(f"⚠️ No text extracted from {file.filename}, using placeholder text")

        logger.info(f"📄 Extracted {len(full_text)} characters from {file.filename}")

        # ---- chunking (UNCHANGED LOGIC) ----
        sentences = re.split(r'(?<=[.!?])\s+', full_text)
        chunks, current = [], ""

        for s in sentences:
            if len(current) + len(s) < 800:
                current += s + " "
            else:
                chunks.append(current.strip())
                current = s + " "

        if current.strip():
            chunks.append(current.strip())

        logger.info(f"🧩 Chunks created: {len(chunks)}")
        logger.info(f"📏 First chunk length: {len(chunks[0]) if chunks else 0}")


        # ---- pinecone upsert (UNCHANGED) ----
        logger.info(f"🔍 Pinecone check: pine_index={pine_index is not None}, chunks={len(chunks)}")
        
        if pine_index:
            vectors = []
            ts = int(time.time())

            for i, chunk in enumerate(chunks):
                vectors.append({
                    "id": f"{user_id}:{file.filename}:{i}:{ts}",
                    "values": get_embedding(chunk),
                    "metadata": {
                        "text": chunk,
                        "user_id": user_id,
                        "source": file.filename,
                        "chunk": i,
                        "uploaded_at": datetime.utcnow().isoformat()
                    }
                })

            logger.info(f"📤 Upserting {len(vectors)} vectors to Pinecone")
            for i in range(0, len(vectors), 100):
                pine_index.upsert(vectors=vectors[i:i+100])
            logger.info(f"✅ Successfully upserted vectors to Pinecone")
        else:
            logger.warning("⚠️ Pinecone index not available - skipping vector storage")

        return {
            "success": True,
            "filename": file.filename,
            "chunks": len(chunks),
            "message": f"{file.filename} uploaded and indexed successfully"
        }

    except Exception as e:
        logger.error(traceback.format_exc())
        raise HTTPException(500, f"Failed to process file: {str(e)}")


@app.get("/api/files")
async def list_files(request: Request):
    user_id = get_current_user_id(request)
    user_dir = os.path.join(UPLOAD_DIR, user_id)

    if not os.path.exists(user_dir):
        return {"success": True, "data": {"files": []}}

    files = []
    for name in os.listdir(user_dir):
        path = os.path.join(user_dir, name)
        if os.path.isfile(path):
            stat = os.stat(path)
            files.append({
                "id": name,
                "name": name,
                "type": "PDF",
                "size": f"{stat.st_size // 1024} KB",
                "uploadDate": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d")
            })

    return {
        "success": True,
        "data": {"files": files}
    }

@app.post("/rag/delete")
async def delete_file(request: Request, delete_request: DeleteFileRequest):
    """Delete file from local storage and Pinecone"""
    user_id = get_current_user_id(request)
    filename = delete_request.filename
    
    try:
        # Delete from local storage
        user_dir = os.path.join(UPLOAD_DIR, user_id)
        file_path = os.path.join(user_dir, filename)
        
        if os.path.exists(file_path):
            os.remove(file_path)
            logger.info(f"Deleted local file: {file_path}")
        
        # Try to delete from Pinecone (optional - don't fail if this fails)
        try:
            if PINECONE_API_KEY:
                pc = Pinecone(api_key=PINECONE_API_KEY)
                index = pc.Index(PINECONE_INDEX)
                
                # Delete vectors with metadata matching the source and user_id
                index.delete(filter={"source": filename, "user_id": user_id})
                logger.info(f"Deleted from Pinecone: {filename} for user {user_id}")
        except Exception as pinecone_error:
            logger.warning(f"Failed to delete from Pinecone: {pinecone_error}")
            # Continue anyway - local file deletion is more important
        
        return {"success": True, "message": f"File {filename} deleted successfully"}
    
    except Exception as e:
        logger.error(f"Error deleting file {filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to delete file: {str(e)}")

@app.get("/uploads/{user_id}/{filename}")
async def serve_file(user_id: str, filename: str):
    """Serve uploaded files - public access with user_id in path"""
    user_dir = os.path.join(UPLOAD_DIR, user_id)
    file_path = os.path.join(user_dir, filename)
    
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    return FileResponse(
        file_path, 
        media_type='application/pdf',
        headers={"Content-Disposition": "inline"}
    )

@app.get("/health")
async def health_check():
    """
    Health check endpoint to verify service status.
    
    Returns:
        JSON with service status information
    """
    collections = []
    mongo_status = "disconnected"
    bedrock_status = "disconnected"
    pinecone_status = "disconnected"
    
    if db is not None:
        try:
            collections = db.list_collection_names()
            mongo_status = "connected"
        except Exception as e:
            logger.error(f"MongoDB health check failed: {e}")
    
    if bedrock:
        bedrock_status = "connected"
    
    if pine_index:
        pinecone_status = "connected"
    
    return {
        "status": "running",
        "service": "Hybrid RAG Assistant",
        "version": "4.0.0",
        "timestamp": datetime.now().isoformat(),
        "database": mongo_status,
        "collections": collections,
        "bedrock": bedrock_status,
        "pinecone": pinecone_status
    }

@app.get("/")
async def root():
    """
    Root endpoint with API information.
    
    Returns:
        JSON with API information
    """
    return {
        "message": "Hybrid RAG Assistant API",
        "version": "4.0.0",
        "description": "AI Assistant with enhanced response formatting (Style A)",
        "endpoints": [
            "POST /api/chat - Send a message (form fields: 'query', 'user_id')",
            "POST /api/upload-pdf - Upload PDF document (form fields: 'file', 'user_id')",
            "GET /health - System status",
            "GET / - This information page"
        ],
        "documentation": "See /docs for interactive API documentation"
    }

# ============================================================
# Main
# ============================================================
if __name__ == "__main__":
    print("Starting Hybrid RAG Assistant with Enhanced Formatting (Style A)")
    print("=" * 50)
    print(f"Model: {BEDROCK_MODEL_ID}")
    print(f"Database: {MONGO_DB}")
    print(f"Pinecone: {PINECONE_INDEX}")
    print("=" * 50)
    print("API Endpoints:")
    print("  POST /api/chat")
    print("  POST /api/upload-pdf")
    print("  GET  /health")
    print("  GET  /")
    print("=" * 50)
    print(f"Server running on http://0.0.0.0:8000")
    print("Press Ctrl+C to stop")
    print("=" * 50)
    
    uvicorn.run(
        app,
        host="0.0.0.0",
        port=int(os.getenv("PORTT", 8000)),
        timeout_keep_alive=60,
        access_log=True,
    )